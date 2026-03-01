#!/usr/bin/env python3
"""
OvaSearch Document Preprocessor
Extracts text from PDF, DOCX, and PPTX files into plain text for ingestion.
Usage: python3 prepare_documents.py [data_folder]
"""
import os
import sys
import importlib

DATA_DIR = sys.argv[1] if len(sys.argv) > 1 else os.path.join(os.path.dirname(__file__), "data")

def check_dependencies():
    missing = []
    for pkg, pip_name in [("pdfplumber", "pdfplumber"), ("docx", "python-docx"), ("pptx", "python-pptx")]:
        try:
            importlib.import_module(pkg)
        except ImportError:
            missing.append(pip_name)
    if missing:
        print(f"Installing: {', '.join(missing)}")
        os.system(f"{sys.executable} -m pip install -q {' '.join(missing)}")

def extract_pdf(filepath):
    import pdfplumber
    text = []
    with pdfplumber.open(filepath) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text.append(t)
    return "\n\n".join(text)

def extract_docx(filepath):
    from docx import Document
    doc = Document(filepath)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

def extract_pptx(filepath):
    from pptx import Presentation
    prs = Presentation(filepath)
    text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if slide_text:
            text.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
    return "\n\n".join(text)

EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
}

def main():
    if not os.path.isdir(DATA_DIR):
        print(f"Data directory not found: {DATA_DIR}")
        sys.exit(1)

    check_dependencies()

    processed = 0
    for root, _, files in os.walk(DATA_DIR):
        for fname in files:
            ext = os.path.splitext(fname)[1].lower()
            if ext not in EXTRACTORS:
                continue

            filepath = os.path.join(root, fname)
            outpath = filepath + ".extracted.txt"

            if os.path.exists(outpath) and os.path.getmtime(outpath) > os.path.getmtime(filepath):
                print(f"  [skip] {fname} (already extracted)")
                continue

            print(f"  [extract] {fname} -> {fname}.extracted.txt")
            try:
                text = EXTRACTORS[ext](filepath)
                with open(outpath, "w", encoding="utf-8") as f:
                    f.write(text)
                processed += 1
            except Exception as e:
                print(f"  [error] {fname}: {e}")

    if processed == 0:
        print("No new documents to process.")
    else:
        print(f"\nExtracted {processed} document(s). Run OvaSearch to index them.")

if __name__ == "__main__":
    main()
